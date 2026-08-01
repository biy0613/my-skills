"""
pptx_helpers.py
PPTX Design System v1.0 — python-pptx 헬퍼 번들

이 모듈은 pptx-design-system 스킬의 핵심 실행 도구다.
새 PPT 만들 때 매번 이 함수들을 다시 짜지 말 것.
모든 컴포넌트는 이 모듈의 함수를 호출해서 일관성 유지.

Usage:
    from scripts.pptx_helpers import *
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    add_section_header(slide, "Stage 1", subtitle="Day 1 매출 검증")
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================
# Color Tokens
# ============================================================

# Surface
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xFA, 0xF6, 0xEE)
PEACH = RGBColor(0xFC, 0xEF, 0xE6)
LINE = RGBColor(0xE8, 0xE1, 0xD4)
LINE_SOFT = RGBColor(0xF0, 0xEB, 0xE0)

# Ink
INK = RGBColor(0x16, 0x22, 0x3F)
INK_SOFT = RGBColor(0x3B, 0x48, 0x6A)
MUTE = RGBColor(0x8C, 0x84, 0x78)

# Accent (coral)
ACCENT = RGBColor(0xE6, 0x5A, 0x3A)
ACCENT_DEEP = RGBColor(0xC9, 0x45, 0x24)
ACCENT_SOFT = RGBColor(0xFD, 0xE6, 0xDA)

# Semantic signal colors
SIGNAL_GREEN = RGBColor(0x3A, 0x8F, 0x5C)
SIGNAL_YELLOW = RGBColor(0xC4, 0x9B, 0x2F)
SIGNAL_RED = ACCENT_DEEP  # alias

# Aliases for emphasis (semantics)
INK_BOLD = INK
ACCENT_BOLD = ACCENT

# ============================================================
# Font / Size Defaults
# ============================================================

FONT_PRIMARY = 'Pretendard'
FONT_FALLBACK = '맑은 고딕'

# Standard sizes (Pt)
SIZE_DISPLAY = 40
SIZE_H1 = 30
SIZE_H2 = 23
SIZE_H3 = 19
SIZE_BODY_LG = 17
SIZE_BODY = 15
SIZE_LABEL = 12
SIZE_CAPTION = 11
SIZE_OVERLINE = 10

# Standard slide dims
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Standard margins
MARGIN_LEFT = Inches(0.8)
MARGIN_RIGHT = Inches(0.5)
MARGIN_TOP = Inches(0.5)
MARGIN_BOTTOM = Inches(0.4)
CONTENT_W = Inches(12.0)


# ============================================================
# Foundation: Background, Presentation Setup
# ============================================================

def new_presentation():
    """새 16:9 presentation 생성. 표준 dim 설정 후 반환."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank_slide(prs):
    """Blank layout 슬라이드 추가. 배경 white로 자동 설정."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    return slide


def set_slide_bg(slide):
    """슬라이드 배경 white. 다른 색 절대 금지."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


# ============================================================
# Text Primitives
# ============================================================

def add_text_box(
    slide, left, top, width, height, text,
    font_size=SIZE_BODY, bold=False, color=INK,
    font_name=FONT_PRIMARY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
):
    """
    단순 텍스트 박스. 단일 스타일.
    Returns: shape (textbox)
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_rich_text(
    tf_or_shape, segments,
    font_size=SIZE_BODY, font_name=FONT_PRIMARY, align=PP_ALIGN.LEFT,
):
    """
    같은 paragraph 안에서 run을 분리해 인라인 강조 적용.

    segments: list of (text, bold, color) tuples
        예: [
            ("결론은 ", False, INK),
            ("고점권 선별", True, ACCENT),  # accent-bold
            ("이다.", False, INK),
        ]

    tf_or_shape: text_frame 또는 shape (shape면 .text_frame 자동 추출)
    """
    tf = tf_or_shape.text_frame if hasattr(tf_or_shape, 'text_frame') else tf_or_shape
    tf.word_wrap = True

    # 첫 paragraph 재사용 (clear)
    p = tf.paragraphs[0]
    p.alignment = align
    # 기존 run 제거
    for r in list(p.runs):
        r._r.getparent().remove(r._r)

    for text, bold, color in segments:
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


# ============================================================
# Section Header (틱마크 + 타이틀 + 부제)
# ============================================================

def add_tick_mark(slide, left=None, top=None, height=Inches(0.45)):
    """좌측 코랄 4px (≈Pt 5) 틱마크. 보더 없음."""
    if left is None:
        left = MARGIN_LEFT
    if top is None:
        top = MARGIN_TOP
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Pt(5), height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()  # 보더 제거
    return bar


def add_section_header(
    slide, title, subtitle=None,
    left=MARGIN_LEFT, top=MARGIN_TOP,
    title_size=SIZE_H1, subtitle_size=SIZE_BODY_LG, subtitle_color=None,
    title_offset=Inches(0.3),
):
    """
    슬라이드 상단 표준 헤더:
    [틱마크] [타이틀 ink bold] / [부제 ink-soft 또는 accent]
    """
    if subtitle_color is None:
        subtitle_color = INK_SOFT

    # Tick mark
    add_tick_mark(slide, left=left, top=top, height=Inches(0.45))

    # Title
    add_text_box(
        slide,
        left + title_offset, top - Inches(0.05),
        CONTENT_W, Inches(0.7),
        title, font_size=title_size, bold=True, color=INK,
    )

    # Subtitle
    if subtitle:
        add_text_box(
            slide,
            left + title_offset, top + Inches(0.55),
            CONTENT_W, Inches(0.5),
            subtitle, font_size=subtitle_size, bold=False, color=subtitle_color,
        )


# ============================================================
# Card (cream background)
# ============================================================

def add_card(slide, left, top, width, height, fill=None, border=None):
    """
    Cream 카드. ROUNDED_RECTANGLE + LINE 1pt 보더.
    fill/border override 가능 (예: PEACH 카드, 보더 제거).
    Returns: shape
    """
    if fill is None:
        fill = CREAM
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border is None:
        shape.line.color.rgb = LINE
        shape.line.width = Pt(1)
    elif border is False:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    # 텍스트는 나중에 add_text_box로 카드 위에 별도 배치 (여백 제어 용이)
    return shape


# ============================================================
# Callout (peach + left accent bar)
# ============================================================

def add_callout(
    slide, left, top, width, height,
    label="Key Takeaway", body=None,
    label_color=None, body_color=None,
):
    """
    Peach 배경 + 좌측 코랄 accent bar.
    핵심 takeaway, 경고, 결론 박스에 사용.

    label: 상단 라벨 (UPPER 권장, accent-deep)
    body: 본문 문자열 또는 segments list (rich text)
    """
    if label_color is None:
        label_color = ACCENT_DEEP
    if body_color is None:
        body_color = INK

    # Peach background (rounded)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PEACH
    bg.line.fill.background()

    # Left accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Pt(5), height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Label
    pad_left = left + Inches(0.25)
    text_w = width - Inches(0.4)
    add_text_box(
        slide,
        pad_left, top + Inches(0.15),
        text_w, Inches(0.35),
        label.upper(),
        font_size=SIZE_OVERLINE, bold=True, color=label_color,
    )

    # Body
    if body is not None:
        body_top = top + Inches(0.5)
        body_h = height - Inches(0.6)
        if isinstance(body, list):
            # rich text segments
            tb = slide.shapes.add_textbox(pad_left, body_top, text_w, body_h)
            tb.text_frame.word_wrap = True
            add_rich_text(tb, body, font_size=SIZE_BODY_LG)
        else:
            add_text_box(
                slide, pad_left, body_top, text_w, body_h,
                body, font_size=SIZE_BODY_LG, bold=False, color=body_color,
            )

    return bg


# ============================================================
# Table styling
# ============================================================

def style_table(table, header_row=True, header_size=SIZE_LABEL, body_size=13):
    """
    Editorial Research 테이블 스타일.
    헤더: cream bg + ink bold
    본문: white bg + ink
    """
    n_rows = len(table.rows)
    n_cols = len(table.columns)

    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)

            if header_row and r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CREAM
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.bold = True
                        run.font.color.rgb = INK
                        run.font.size = Pt(header_size)
                        run.font.name = FONT_PRIMARY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.color.rgb = INK
                        run.font.size = Pt(body_size)
                        run.font.name = FONT_PRIMARY


def add_styled_table(slide, rows_data, left, top, width, height, header_row=True):
    """
    rows_data: list[list[str]] — 첫 행이 헤더(header_row=True 시)
    한 번에 테이블 생성 + 데이터 채우기 + 스타일링.
    """
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if n_rows else 0
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
    style_table(table, header_row=header_row)
    return table


# ============================================================
# Chip (priority/status badge)
# ============================================================

def add_chip(
    slide, left, top, text,
    variant='accent', width=None, height=None,
):
    """
    Pill-style chip.
    variant: 'accent' (코랄 bg + white) | 'ink' (네이비 bg + white)
        | 'mute' (mute bg + white) | 'green' | 'yellow'
    """
    if width is None:
        # 글자 길이 기반 가변 (대략)
        width = Inches(max(0.4, 0.18 * max(len(text), 2)))
    if height is None:
        height = Inches(0.32)

    color_map = {
        'accent': ACCENT,
        'ink': INK,
        'mute': MUTE,
        'green': SIGNAL_GREEN,
        'yellow': SIGNAL_YELLOW,
    }
    bg_color = color_map.get(variant, ACCENT)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_PRIMARY
    run.font.size = Pt(SIZE_CAPTION)
    run.font.bold = True
    run.font.color.rgb = WHITE
    return shape


# ============================================================
# Bullet helpers (signal lists)
# ============================================================

def add_bullet_list(
    slide, left, top, width, height, items,
    dot_color=ACCENT, font_size=SIZE_BODY,
):
    """
    items: list of strings 또는 list of segments-list (rich text).
    각 item 앞에 컬러 도트(•) 추가.
    Yellow flag → dot_color=ACCENT_DEEP
    Good signal → dot_color=SIGNAL_GREEN
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)

        # Dot
        dot = p.add_run()
        dot.text = "● "
        dot.font.name = FONT_PRIMARY
        dot.font.size = Pt(font_size)
        dot.font.color.rgb = dot_color

        # Content
        if isinstance(item, list):
            # rich segments
            for text, bold, color in item:
                run = p.add_run()
                run.text = text
                run.font.name = FONT_PRIMARY
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = str(item)
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(font_size)
            run.font.color.rgb = INK
    return tb


# ============================================================
# Convenience: Title slide
# ============================================================

def build_title_slide(prs, title, subtitle=None, meta=None):
    """
    표준 Title 슬라이드 한 줄 생성.
    title: Display 40pt ink bold
    subtitle: 18pt accent
    meta: "발표자 | 소속 | 날짜" — 14pt mute
    """
    slide = add_blank_slide(prs)
    cx = MARGIN_LEFT
    cy = Inches(2.5)

    add_text_box(
        slide, cx, cy, CONTENT_W, Inches(1.0),
        title, font_size=SIZE_DISPLAY, bold=True, color=INK,
    )
    cy += Inches(1.1)

    if subtitle:
        add_text_box(
            slide, cx, cy, CONTENT_W, Inches(0.6),
            subtitle, font_size=SIZE_BODY_LG, bold=False, color=ACCENT,
        )
        cy += Inches(0.7)

    if meta:
        add_text_box(
            slide, cx, cy, CONTENT_W, Inches(0.4),
            meta, font_size=14, bold=False, color=MUTE,
        )

    return slide


# ============================================================
# Convenience: LIVE DEMO badge
# ============================================================

def add_live_demo_badge(slide, top=None, right_margin=Inches(0.5)):
    """우상단 LIVE DEMO 칩."""
    if top is None:
        top = MARGIN_TOP
    width = Inches(1.2)
    left = SLIDE_W - right_margin - width
    return add_chip(
        slide, left, top, "LIVE DEMO",
        variant='accent', width=width, height=Inches(0.35),
    )
